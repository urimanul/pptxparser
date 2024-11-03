import streamlit as st
from pptx import Presentation
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lex_rank import LexRankSummarizer
from io import BytesIO
import mysql.connector

def read_pptx(file):
    pptx_reader = Presentation(file)
    text = ""
    for slide in pptx_reader.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def save_summary_to_db(uploaded_file, summary):
    # MySQLデータベースに接続
    conn = mysql.connector.connect(
        host="www.ryhintl.com",
        user="smairuser",
        password="smairuser",
        database="smair",
        port=36000
    )
    cursor = conn.cursor()

    # 要約結果をデータベースに挿入
    cursor.execute("INSERT INTO parsed_summary (uploaded_file, summary_text) VALUES (%s, %s)", (uploaded_file.name, summary))
    conn.commit()

    cursor.close()
    conn.close()

# Streamlitアプリの設定
st.title('PPTXパーサー')
st.write('PPTXファイル要約')

uploaded_file = st.file_uploader("PPTXファイルを選択してください", type=["pptx"])

if uploaded_file is not None:
    text = read_pptx(uploaded_file)
    
    # 読み込んだPPTXファイルの文字数を表示
    st.write(f"読み込んだPPTXファイルの文字数: {len(text)}文字")
    
    # 要約文字数を指定
    char_count = st.number_input("要約文字数を指定してください", min_value=1, max_value=100000, value=2000)
    
    # 要約を実行
    parser = PlaintextParser.from_string(text, Tokenizer('japanese'))
    summarizer = LexRankSummarizer()
    res = summarizer(document=parser.document, sentences_count=char_count)  # 大きな値を設定しておく

    # 指定された文字数に要約を調整
    summary = ""
    for sentence in res:
        if len(summary) + len(str(sentence)) <= char_count:
            summary += str(sentence)
        else:
            break

    # フォントサイズを12pxに設定して表示
    st.markdown(f"<div style='font-size: 12px;'>{summary}</div>", unsafe_allow_html=True)

    # 登録ボタンを追加
    if st.button("登録"):
        save_summary_to_db(uploaded_file, summary)
        st.success("要約結果がデータベースに登録されました。")
